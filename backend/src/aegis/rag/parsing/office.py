"""DOCX, PPTX, and XLSX parsing.

Office formats are the easy case structurally — they carry explicit style information — and
the hard case in terms of what a naive extractor loses:

* **DOCX** keeps paragraphs in ``document.paragraphs`` and tables in ``document.tables``, as
  two separate lists. Iterating them in that order relocates every table to the end of the
  document, so the underlying XML body is walked instead to preserve reading order.
* **PPTX** has no reading order at all. Shapes are stored in z-order, so a deck extracts as
  scrambled fragments unless shapes are sorted by position. Speaker notes carry most of the
  actual explanation in a corporate deck and are extracted separately.
* **XLSX** formulas are worthless to a retriever (``=SUM(B2:B40)`` answers nothing) while
  cached values are exactly what a question is about, so cells are read with
  ``data_only=True``. Merged headers are unmerged so every row keeps its column names.
"""

from __future__ import annotations

import asyncio
import io
from typing import Any

from aegis.core.errors import ParserError
from aegis.core.logging import get_logger
from aegis.domain.enums import ChunkType
from aegis.domain.ports.parser import DocumentBlock, PageInfo, ParsedDocument
from aegis.rag.parsing.cleaning import clean_text, is_meaningful, is_noise, normalize_whitespace
from aegis.rag.parsing.sniff import DOCX, PPTX, XLSX
from aegis.rag.parsing.structure import assign_hierarchy, document_title
from aegis.rag.parsing.tables import looks_like_header, to_markdown, to_row_records

logger = get_logger(__name__)

MAX_SHEET_ROWS = 20_000
_LARGE_TABLE_ROWS = 12
# Word's built-in heading styles map directly to levels; "Title" is treated as level 1
# because that is how authors use it.
_WORD_HEADINGS = {
    "title": 1,
    "heading 1": 1,
    "heading 2": 2,
    "heading 3": 3,
    "heading 4": 4,
    "heading 5": 5,
    "heading 6": 6,
    "subtitle": 2,
}


class DocxParser:
    """Word documents, walked in body order so tables stay where the author put them."""

    name = "docx"

    def supports(self, mime_type: str, filename: str) -> bool:
        return mime_type == DOCX

    async def parse(self, data: bytes, *, filename: str, mime_type: str) -> ParsedDocument:
        return await asyncio.to_thread(self._parse_sync, data, filename)

    def _parse_sync(self, data: bytes, filename: str) -> ParsedDocument:
        from docx import Document
        from docx.table import Table
        from docx.text.paragraph import Paragraph

        try:
            document = Document(io.BytesIO(data))
        except Exception as exc:
            raise ParserError(f"The Word document could not be read: {exc}") from exc

        blocks: list[DocumentBlock] = []
        offset = 0

        for element in _iter_body(document):
            if isinstance(element, Paragraph):
                block = self._paragraph_block(element, offset)
                if block is not None:
                    offset = block.char_end
                    blocks.append(block)
            elif isinstance(element, Table):
                blocks.extend(self._table_blocks(element))

        # Headers and footers are extracted once, not per page: python-docx has no concept of
        # pages, and the content is identical on each.
        blocks.extend(self._section_extras(document))

        core = document.core_properties
        assign_hierarchy(blocks)
        return ParsedDocument(
            blocks=blocks,
            pages=[PageInfo(number=1, char_count=sum(len(b.text) for b in blocks))],
            title=(core.title or "").strip() or document_title(blocks, fallback=filename),
            parser=self.name,
            metadata={
                key: str(value)[:500]
                for key, value in (
                    ("author", core.author),
                    ("subject", core.subject),
                    ("category", core.category),
                    ("last_modified_by", core.last_modified_by),
                )
                if value
            },
        )

    def _paragraph_block(self, paragraph: Any, offset: int) -> DocumentBlock | None:
        text = clean_text(paragraph.text)
        if not text:
            return None

        style = (paragraph.style.name or "").strip().lower() if paragraph.style else ""
        level = _WORD_HEADINGS.get(style)
        if level is None and style.startswith("heading"):
            # "Heading 7"+ and localised variants: keep the content, clamp the depth.
            level = 6

        if level is not None:
            return DocumentBlock(
                text=normalize_whitespace(text),
                block_type=ChunkType.HEADING,
                heading_level=level,
                char_start=offset,
                char_end=offset + len(text),
                metadata={"style": style},
            )

        block_type = ChunkType.TEXT
        if style.startswith(("list", "bullet")):
            block_type = ChunkType.LIST
            text = f"- {text}"
        elif "caption" in style:
            block_type = ChunkType.CAPTION
        elif "quote" in style:
            block_type = ChunkType.TEXT

        if block_type is ChunkType.TEXT and is_noise(text):
            return None
        return DocumentBlock(
            text=text,
            block_type=block_type,
            char_start=offset,
            char_end=offset + len(text),
        )

    def _table_blocks(self, table: Any) -> list[DocumentBlock]:
        rows = [
            [clean_text(cell.text, preserve_layout=False) for cell in row.cells]
            for row in table.rows
        ]
        rows = [row for row in rows if any(row)]
        if len(rows) < 2:
            # A single-row "table" is almost always a layout device, not data.
            return (
                [DocumentBlock(text=" ".join(rows[0]), block_type=ChunkType.TEXT)]
                if rows and is_meaningful(" ".join(rows[0]))
                else []
            )

        header = rows[0] if looks_like_header(rows[0]) else []
        body = rows[1:] if header else rows
        columns = header or [f"column_{i + 1}" for i in range(len(rows[0]))]

        if len(rows) <= _LARGE_TABLE_ROWS:
            return [
                DocumentBlock(
                    text=to_markdown(body, header=columns),
                    block_type=ChunkType.TABLE,
                    metadata={"format": "markdown", "rows": len(body)},
                )
            ]

        records = to_row_records(body, header=columns)
        return [
            DocumentBlock(
                text="\n".join(records[start : start + 20]),
                block_type=ChunkType.TABLE,
                metadata={
                    "format": "records",
                    "row_from": start + 1,
                    "row_to": min(start + 20, len(records)),
                    "columns": [str(c) for c in columns],
                },
            )
            for start in range(0, len(records), 20)
        ]

    def _section_extras(self, document: Any) -> list[DocumentBlock]:
        """Headers and footers, deduplicated.

        Included because a document's header often carries its only classification marking
        ("Confidential — Internal Use"), which is exactly the kind of thing a reviewer needs
        to see in a citation.
        """
        seen: set[str] = set()
        blocks: list[DocumentBlock] = []
        for section in document.sections:
            for part, label in ((section.header, "header"), (section.footer, "footer")):
                if part is None:
                    continue
                text = clean_text("\n".join(p.text for p in part.paragraphs))
                if not text or text in seen or len(text) < 8:
                    continue
                seen.add(text)
                blocks.append(
                    DocumentBlock(
                        text=text,
                        block_type=ChunkType.CAPTION,
                        metadata={"role": label},
                    )
                )
        return blocks


class PptxParser:
    """PowerPoint decks. One page per slide, with notes attached.

    Reading order is reconstructed from shape geometry: PowerPoint stores shapes in z-order,
    so a title added last appears last in the XML. Sorting by top-then-left position recovers
    what a human sees, and without it a deck extracts as an unreadable shuffle.
    """

    name = "pptx"

    def supports(self, mime_type: str, filename: str) -> bool:
        return mime_type == PPTX

    async def parse(self, data: bytes, *, filename: str, mime_type: str) -> ParsedDocument:
        return await asyncio.to_thread(self._parse_sync, data, filename)

    def _parse_sync(self, data: bytes, filename: str) -> ParsedDocument:
        from pptx import Presentation

        try:
            deck = Presentation(io.BytesIO(data))
        except Exception as exc:
            raise ParserError(f"The presentation could not be read: {exc}") from exc

        blocks: list[DocumentBlock] = []
        pages: list[PageInfo] = []

        for number, slide in enumerate(deck.slides, start=1):
            slide_blocks = self._slide_blocks(slide, number)
            blocks.extend(slide_blocks)
            pages.append(PageInfo(number=number, char_count=sum(len(b.text) for b in slide_blocks)))

        assign_hierarchy(blocks)
        core = deck.core_properties
        return ParsedDocument(
            blocks=blocks,
            pages=pages,
            title=(core.title or "").strip() or document_title(blocks, fallback=filename),
            parser=self.name,
            metadata={"slide_count": len(pages)},
        )

    def _slide_blocks(self, slide: Any, number: int) -> list[DocumentBlock]:
        blocks: list[DocumentBlock] = []
        title_text = ""
        if slide.shapes.title is not None and slide.shapes.title.has_text_frame:
            title_text = clean_text(slide.shapes.title.text)
            if title_text:
                blocks.append(
                    DocumentBlock(
                        text=normalize_whitespace(title_text),
                        block_type=ChunkType.HEADING,
                        page=number,
                        # Slide titles are siblings, not a hierarchy: nesting slide 4 under
                        # slide 3 would invent a relationship the deck does not have.
                        heading_level=1,
                        metadata={"slide": number},
                    )
                )

        ordered = sorted(
            slide.shapes,
            key=lambda shape: (
                float(getattr(shape, "top", 0) or 0),
                float(getattr(shape, "left", 0) or 0),
            ),
        )
        body: list[str] = []
        for shape in ordered:
            if shape is slide.shapes.title:
                continue
            if getattr(shape, "has_table", False):
                rows = [[clean_text(cell.text) for cell in row.cells] for row in shape.table.rows]
                rows = [r for r in rows if any(r)]
                if len(rows) >= 2:
                    blocks.append(
                        DocumentBlock(
                            text=to_markdown(rows),
                            block_type=ChunkType.TABLE,
                            page=number,
                            metadata={"format": "markdown", "slide": number},
                        )
                    )
                continue
            if getattr(shape, "has_text_frame", False):
                text = clean_text(shape.text_frame.text)
                if text and text != title_text:
                    body.append(text)

        if body:
            blocks.append(
                DocumentBlock(
                    text="\n".join(body),
                    block_type=ChunkType.TEXT,
                    page=number,
                    metadata={"slide": number},
                )
            )

        notes = self._notes(slide)
        if notes:
            # Notes are kept as a separate block, not merged into the slide body: they are the
            # narrative a slide's bullet points omit, and they answer questions the slide
            # itself cannot.
            blocks.append(
                DocumentBlock(
                    text=notes,
                    block_type=ChunkType.TEXT,
                    page=number,
                    metadata={"slide": number, "role": "speaker_notes"},
                )
            )
        return blocks

    def _notes(self, slide: Any) -> str:
        if not getattr(slide, "has_notes_slide", False):
            return ""
        frame = getattr(slide.notes_slide, "notes_text_frame", None)
        if frame is None:
            return ""
        text = clean_text(frame.text)
        return text if is_meaningful(text) else ""


class XlsxParser:
    """Excel workbooks. One page per sheet, rows as records.

    ``data_only=True`` reads cached formula results. The alternative — formula strings — is
    both unretrievable and misleading: nobody asks a question whose answer is
    ``=VLOOKUP(A2,Rates!A:B,2,FALSE)``.
    """

    name = "xlsx"

    def supports(self, mime_type: str, filename: str) -> bool:
        return mime_type == XLSX

    async def parse(self, data: bytes, *, filename: str, mime_type: str) -> ParsedDocument:
        return await asyncio.to_thread(self._parse_sync, data, filename)

    def _parse_sync(self, data: bytes, filename: str) -> ParsedDocument:
        from openpyxl import load_workbook

        try:
            workbook = load_workbook(
                io.BytesIO(data), data_only=True, read_only=True, keep_links=False
            )
        except Exception as exc:
            raise ParserError(f"The workbook could not be read: {exc}") from exc

        blocks: list[DocumentBlock] = []
        pages: list[PageInfo] = []
        warnings: list[str] = []

        try:
            for number, sheet in enumerate(workbook.worksheets, start=1):
                sheet_blocks, truncated = self._sheet_blocks(sheet, number)
                if truncated:
                    warnings.append(
                        f"Sheet {sheet.title!r}: only the first "
                        f"{MAX_SHEET_ROWS:,} rows were indexed."
                    )
                blocks.extend(sheet_blocks)
                pages.append(
                    PageInfo(number=number, char_count=sum(len(b.text) for b in sheet_blocks))
                )
        finally:
            workbook.close()

        assign_hierarchy(blocks)
        return ParsedDocument(
            blocks=blocks,
            pages=pages,
            title=filename,
            parser=self.name,
            metadata={"sheet_count": len(pages)},
            warnings=warnings,
        )

    def _sheet_blocks(self, sheet: Any, number: int) -> tuple[list[DocumentBlock], bool]:
        rows: list[list[str]] = []
        truncated = False
        for index, row in enumerate(sheet.iter_rows(values_only=True)):
            if index >= MAX_SHEET_ROWS:
                truncated = True
                break
            cells = ["" if value is None else str(value).strip() for value in row]
            if any(cells):
                rows.append(cells)

        if not rows:
            return [], truncated

        # The sheet name is a heading so that "which sheet did this come from" is answerable
        # from the citation alone.
        blocks: list[DocumentBlock] = [
            DocumentBlock(
                text=str(sheet.title),
                block_type=ChunkType.HEADING,
                page=number,
                heading_level=1,
                metadata={"sheet": str(sheet.title)},
            )
        ]

        header = rows[0] if looks_like_header(rows[0]) else []
        body = rows[1:] if header else rows
        columns = header or [f"column_{i + 1}" for i in range(max(len(r) for r in rows))]

        records = to_row_records(body, header=columns)
        group = 25
        for start in range(0, len(records), group):
            window = records[start : start + group]
            blocks.append(
                DocumentBlock(
                    text="\n".join(window),
                    block_type=ChunkType.TABLE,
                    page=number,
                    metadata={
                        "format": "records",
                        "sheet": str(sheet.title),
                        "row_from": start + 1,
                        "row_to": start + len(window),
                        "columns": [str(c) for c in columns],
                    },
                )
            )
        return blocks, truncated


def _iter_body(document: Any) -> list[Any]:
    """Yield paragraphs and tables in the order they appear in the document body.

    ``python-docx`` exposes them as two separate collections, which loses interleaving. The
    body XML is the only place the real order exists, so it is walked directly — the one place
    in this module that reaches past the library's public surface, and the alternative is
    every table in the document moving to the end.
    """
    from docx.oxml.ns import qn
    from docx.table import Table
    from docx.text.paragraph import Paragraph

    body = document.element.body
    items: list[Any] = []
    for child in body.iterchildren():
        if child.tag == qn("w:p"):
            items.append(Paragraph(child, document))
        elif child.tag == qn("w:tbl"):
            items.append(Table(child, document))
    return items
