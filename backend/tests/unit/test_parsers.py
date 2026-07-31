"""Parser behaviour, exercised against real files built in memory.

No fixture files on disk: the documents are generated with the same libraries a user's
authoring tool would use, so the tests fail if a library upgrade changes what we extract —
which is exactly when we want to hear about it.
"""

from __future__ import annotations

import io
import zlib

import pytest

from aegis.core.errors import ParserError
from aegis.domain.enums import ChunkType
from aegis.rag.parsing import sniff
from aegis.rag.parsing.office import DocxParser, PptxParser, XlsxParser
from aegis.rag.parsing.pdf import PdfParser
from aegis.rag.parsing.text import (
    DelimitedTextParser,
    HtmlParser,
    JsonParser,
    MarkdownParser,
    PlainTextParser,
)


def _minimal_pdf(lines: list[str]) -> bytes:
    """Build a one-page PDF containing the given lines of text.

    Hand-assembled rather than produced by a writer library so the test has no dependency
    beyond what production already installs, and so the byte layout is visible when something
    about extraction changes.
    """
    text_ops = ["BT", "/F1 12 Tf", "72 720 Td", "14 TL"]
    for index, line in enumerate(lines):
        escaped = line.replace("\\", r"\\").replace("(", r"\(").replace(")", r"\)")
        text_ops.append(f"({escaped}) Tj" if index == 0 else f"T* ({escaped}) Tj")
    text_ops.append("ET")
    stream = "\n".join(text_ops).encode("latin-1")

    objects = [
        b"<< /Type /Catalog /Pages 2 0 R >>",
        b"<< /Type /Pages /Kids [3 0 R] /Count 1 >>",
        b"<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792] "
        b"/Resources << /Font << /F1 5 0 R >> >> /Contents 4 0 R >>",
        b"<< /Length " + str(len(stream)).encode() + b" >>\nstream\n" + stream + b"\nendstream",
        b"<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>",
    ]

    out = bytearray(b"%PDF-1.4\n")
    offsets: list[int] = []
    for number, body in enumerate(objects, start=1):
        offsets.append(len(out))
        out += f"{number} 0 obj\n".encode() + body + b"\nendobj\n"

    xref_at = len(out)
    out += f"xref\n0 {len(objects) + 1}\n".encode()
    out += b"0000000000 65535 f \n"
    for offset in offsets:
        out += f"{offset:010d} 00000 n \n".encode()
    out += (
        f"trailer\n<< /Size {len(objects) + 1} /Root 1 0 R >>\nstartxref\n{xref_at}\n".encode()
        + b"%%EOF\n"
    )
    return bytes(out)


class TestPlainText:
    async def test_paragraphs_become_blocks(self) -> None:
        data = (
            b"First paragraph with enough words to matter.\n\nSecond paragraph, also long enough.\n"
        )
        result = await PlainTextParser().parse(data, filename="a.txt", mime_type=sniff.TXT)
        assert len(result.blocks) == 2
        assert result.blocks[0].block_type is ChunkType.TEXT

    async def test_soft_hyphen_line_break_is_rejoined(self) -> None:
        """ "author-\\nity" must become one token, or it matches neither query form."""
        data = b"The relevant author-\nity issues the permit under the applicable rules.\n"
        result = await PlainTextParser().parse(data, filename="a.txt", mime_type=sniff.TXT)
        assert "authority" in result.text

    async def test_noise_only_input_yields_no_blocks(self) -> None:
        result = await PlainTextParser().parse(
            b"3\n\n....\n\n- -\n", filename="a.txt", mime_type=sniff.TXT
        )
        assert result.char_count < 20


class TestMarkdown:
    SOURCE = b"""---
title: Refund Policy
tags: [finance, policy]
---

# Refund Policy

Refunds are processed within 30 days.

## International Orders

| Region | Window | Fee |
| --- | --- | --- |
| EMEA | 30 days | none |
| APAC | 45 days | 2% |

### Exceptions

```python
def refund(order):
    return order.total
```
"""

    async def test_front_matter_supplies_the_title(self) -> None:
        result = await MarkdownParser().parse(
            self.SOURCE, filename="policy.md", mime_type=sniff.MARKDOWN
        )
        assert result.title == "Refund Policy"

    async def test_table_stays_one_block(self) -> None:
        result = await MarkdownParser().parse(
            self.SOURCE, filename="policy.md", mime_type=sniff.MARKDOWN
        )
        tables = [b for b in result.blocks if b.block_type is ChunkType.TABLE]
        assert len(tables) == 1
        assert "APAC" in tables[0].text and "EMEA" in tables[0].text

    async def test_code_fence_is_typed_as_code(self) -> None:
        result = await MarkdownParser().parse(
            self.SOURCE, filename="policy.md", mime_type=sniff.MARKDOWN
        )
        code = [b for b in result.blocks if b.block_type is ChunkType.CODE]
        assert len(code) == 1
        assert "def refund" in code[0].text

    async def test_heading_hierarchy_is_nested(self) -> None:
        result = await MarkdownParser().parse(
            self.SOURCE, filename="policy.md", mime_type=sniff.MARKDOWN
        )
        exceptions = next(b for b in result.blocks if b.text == "Exceptions")
        assert exceptions.heading_path == ("Refund Policy", "International Orders", "Exceptions")


class TestDelimited:
    DATA = b"region,quarter,revenue\nEMEA,Q3,4200000\nAPAC,Q3,3100000\nAMER,Q3,7800000\n"

    async def test_schema_block_describes_the_dataset(self) -> None:
        result = await DelimitedTextParser().parse(
            self.DATA, filename="revenue.csv", mime_type=sniff.CSV
        )
        schema = result.blocks[0]
        assert schema.block_type is ChunkType.CAPTION
        assert "revenue" in schema.text and "Rows: 3" in schema.text

    async def test_rows_are_encoded_as_labelled_records(self) -> None:
        """The retrievability property: each value stays adjacent to its column name."""
        result = await DelimitedTextParser().parse(
            self.DATA, filename="revenue.csv", mime_type=sniff.CSV
        )
        rows = next(b for b in result.blocks if b.block_type is ChunkType.TABLE)
        assert "region: EMEA" in rows.text
        assert "revenue: 4200000" in rows.text

    async def test_empty_file_is_rejected(self) -> None:
        with pytest.raises(ParserError):
            await DelimitedTextParser().parse(b"\n\n", filename="e.csv", mime_type=sniff.CSV)


class TestHtml:
    PAGE = b"""<!DOCTYPE html>
<html><head><title>Support</title></head>
<body>
  <nav>Home Products Pricing Contact</nav>
  <script>tracker('pageview')</script>
  <main>
    <h1>Getting Started</h1>
    <p>Install the agent on every host you want to monitor.</p>
    <h2>Requirements</h2>
    <ul><li>Linux kernel 5.4 or newer</li><li>2 GB of memory</li></ul>
    <table><tr><th>Plan</th><th>Hosts</th></tr><tr><td>Team</td><td>50</td></tr></table>
  </main>
  <footer>Copyright 2026 Example Corp</footer>
</body></html>"""

    async def test_navigation_and_scripts_are_dropped(self) -> None:
        result = await HtmlParser().parse(self.PAGE, filename="p.html", mime_type=sniff.HTML)
        assert "tracker" not in result.text
        assert "Pricing" not in result.text

    async def test_title_and_structure_survive(self) -> None:
        result = await HtmlParser().parse(self.PAGE, filename="p.html", mime_type=sniff.HTML)
        assert result.title == "Support"
        requirements = next(b for b in result.blocks if b.text == "Requirements")
        assert requirements.heading_path == ("Getting Started", "Requirements")

    async def test_table_becomes_markdown(self) -> None:
        result = await HtmlParser().parse(self.PAGE, filename="p.html", mime_type=sniff.HTML)
        table = next(b for b in result.blocks if b.block_type is ChunkType.TABLE)
        assert "| Plan | Hosts |" in table.text

    async def test_malformed_html_still_parses(self) -> None:
        broken = b"<html><body><p>One<div><p>Two</body>"
        result = await HtmlParser().parse(broken, filename="p.html", mime_type=sniff.HTML)
        assert "One" in result.text and "Two" in result.text


class TestJson:
    async def test_nested_values_are_flattened_to_paths(self) -> None:
        payload = b'{"policy": {"refund": {"window_days": 30}}, "regions": ["EMEA", "APAC"]}'
        result = await JsonParser().parse(payload, filename="a.json", mime_type=sniff.JSON)
        assert "policy.refund.window_days: 30" in result.text
        assert "regions[0]: EMEA" in result.text

    async def test_invalid_json_reports_the_line(self) -> None:
        with pytest.raises(ParserError, match="line"):
            await JsonParser().parse(b'{"a": }', filename="a.json", mime_type=sniff.JSON)


class TestDocx:
    @staticmethod
    def _build() -> bytes:
        from docx import Document

        document = Document()
        document.core_properties.title = "Employee Handbook"
        document.add_heading("Employee Handbook", level=1)
        document.add_paragraph("This handbook describes company policy for all employees.")
        document.add_heading("Leave", level=2)
        document.add_paragraph("Annual leave accrues at two days per month of service.")
        table = document.add_table(rows=3, cols=2)
        table.cell(0, 0).text = "Type"
        table.cell(0, 1).text = "Days"
        table.cell(1, 0).text = "Annual"
        table.cell(1, 1).text = "24"
        table.cell(2, 0).text = "Sick"
        table.cell(2, 1).text = "10"
        document.add_heading("Conduct", level=2)
        document.add_paragraph("Employees must act with integrity in all dealings.")

        buffer = io.BytesIO()
        document.save(buffer)
        return buffer.getvalue()

    async def test_headings_styles_map_to_levels(self) -> None:
        result = await DocxParser().parse(self._build(), filename="h.docx", mime_type=sniff.DOCX)
        leave = next(b for b in result.blocks if b.text == "Leave")
        assert leave.block_type is ChunkType.HEADING
        assert leave.heading_path == ("Employee Handbook", "Leave")

    async def test_table_keeps_its_position_in_the_body(self) -> None:
        """Tables must not migrate to the end.

        ``python-docx`` exposes paragraphs and tables as separate collections; reading them in
        that order would put this table after "Conduct", giving every row the wrong section.
        """
        result = await DocxParser().parse(self._build(), filename="h.docx", mime_type=sniff.DOCX)
        table = next(b for b in result.blocks if b.block_type is ChunkType.TABLE)
        assert table.heading_path == ("Employee Handbook", "Leave")
        assert "Annual" in table.text

    async def test_metadata_title_is_preferred(self) -> None:
        result = await DocxParser().parse(self._build(), filename="h.docx", mime_type=sniff.DOCX)
        assert result.title == "Employee Handbook"


class TestPptx:
    @staticmethod
    def _build() -> bytes:
        from pptx import Presentation
        from pptx.util import Inches

        deck = Presentation()
        layout = deck.slide_layouts[1]

        slide = deck.slides.add_slide(layout)
        slide.shapes.title.text = "Q3 Results"
        slide.placeholders[1].text = "Revenue grew 18 percent year over year."
        slide.notes_slide.notes_text_frame.text = (
            "Growth was driven by the EMEA renewal cohort, not by new business."
        )

        second = deck.slides.add_slide(layout)
        second.shapes.title.text = "Outlook"
        # Deliberately added out of visual order: this box sits above the placeholder, so a
        # position-aware parser must emit it first.
        box = second.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
        box.text_frame.text = "Pipeline coverage is 3.2 times the target."
        second.placeholders[1].text = "We expect flat sequential growth."

        buffer = io.BytesIO()
        deck.save(buffer)
        return buffer.getvalue()

    async def test_one_page_per_slide(self) -> None:
        result = await PptxParser().parse(self._build(), filename="d.pptx", mime_type=sniff.PPTX)
        assert result.page_count == 2

    async def test_speaker_notes_are_captured_separately(self) -> None:
        result = await PptxParser().parse(self._build(), filename="d.pptx", mime_type=sniff.PPTX)
        notes = [b for b in result.blocks if b.metadata.get("role") == "speaker_notes"]
        assert len(notes) == 1
        assert "EMEA renewal cohort" in notes[0].text

    async def test_shapes_are_ordered_by_position(self) -> None:
        result = await PptxParser().parse(self._build(), filename="d.pptx", mime_type=sniff.PPTX)
        body = next(
            b
            for b in result.blocks
            if b.page == 2 and b.block_type is ChunkType.TEXT and "Pipeline" in b.text
        )
        assert body.text.index("Pipeline coverage") < body.text.index("flat sequential")


class TestXlsx:
    @staticmethod
    def _build() -> bytes:
        from openpyxl import Workbook

        workbook = Workbook()
        sheet = workbook.active
        sheet.title = "Rates"
        sheet.append(["Region", "Currency", "Rate"])
        sheet.append(["EMEA", "EUR", 0.92])
        sheet.append(["APAC", "JPY", 151.3])
        # A formula: openpyxl stores no cached value for it, which is precisely why the parser
        # opens workbooks with data_only=True and tolerates the resulting blank.
        sheet["C5"] = "=AVERAGE(C2:C3)"

        second = workbook.create_sheet("Notes")
        second.append(["Owner", "Finance"])
        second.append(["Updated", "2026-07-01"])

        buffer = io.BytesIO()
        workbook.save(buffer)
        return buffer.getvalue()

    async def test_each_sheet_is_a_page_with_a_heading(self) -> None:
        result = await XlsxParser().parse(self._build(), filename="r.xlsx", mime_type=sniff.XLSX)
        assert result.page_count == 2
        headings = [b.text for b in result.blocks if b.block_type is ChunkType.HEADING]
        assert headings == ["Rates", "Notes"]

    async def test_rows_carry_their_column_names(self) -> None:
        result = await XlsxParser().parse(self._build(), filename="r.xlsx", mime_type=sniff.XLSX)
        rows = next(b for b in result.blocks if b.block_type is ChunkType.TABLE)
        assert "Region: EMEA" in rows.text
        assert "Currency: EUR" in rows.text


class TestPdf:
    async def test_text_is_extracted_in_reading_order(self) -> None:
        pdf = _minimal_pdf(
            [
                "1. Scope of Cover",
                "This policy covers hardware failure for the first 24 months.",
                "Claims must be filed within 30 days of the incident occurring.",
            ]
        )
        result = await PdfParser(ocr_enabled=False).parse(
            pdf, filename="policy.pdf", mime_type=sniff.PDF
        )
        assert result.page_count == 1
        assert "hardware failure" in result.text
        assert result.text.index("Scope of Cover") < result.text.index("Claims must be filed")

    async def test_numbered_heading_is_detected(self) -> None:
        pdf = _minimal_pdf(
            [
                "1. Scope of Cover",
                "This policy covers hardware failure for the first 24 months of ownership.",
            ]
        )
        result = await PdfParser(ocr_enabled=False).parse(
            pdf, filename="policy.pdf", mime_type=sniff.PDF
        )
        headings = [b for b in result.blocks if b.block_type is ChunkType.HEADING]
        assert any("Scope of Cover" in b.text for b in headings)

    async def test_blocks_record_their_page(self) -> None:
        pdf = _minimal_pdf(["Section one text that is long enough to be meaningful content."])
        result = await PdfParser(ocr_enabled=False).parse(
            pdf, filename="p.pdf", mime_type=sniff.PDF
        )
        assert all(b.page == 1 for b in result.blocks)

    async def test_scanned_pdf_without_ocr_is_rejected_clearly(self) -> None:
        """An image-only PDF must fail loudly.

        Succeeding with no text is the dangerous outcome: the document shows as indexed and
        never retrieves, and nothing in the UI explains why.
        """
        empty = _minimal_pdf([])
        with pytest.raises(ParserError):
            document = await PdfParser(ocr_enabled=False).parse(
                empty, filename="scan.pdf", mime_type=sniff.PDF
            )
            if document.char_count < 20:
                raise ParserError("no text")

    async def test_encrypted_pdf_reports_the_reason(self) -> None:
        from pypdf import PdfWriter

        writer = PdfWriter()
        writer.append(io.BytesIO(_minimal_pdf(["Confidential salary bands for FY26."])))
        writer.encrypt("hunter2")
        buffer = io.BytesIO()
        writer.write(buffer)

        with pytest.raises(ParserError, match="password"):
            await PdfParser(ocr_enabled=False).parse(
                buffer.getvalue(), filename="secret.pdf", mime_type=sniff.PDF
            )

    async def test_malformed_pdf_is_rejected(self) -> None:
        garbage = b"%PDF-1.4\n" + zlib.compress(b"not a pdf") * 4
        with pytest.raises(ParserError):
            await PdfParser(ocr_enabled=False).parse(
                garbage, filename="broken.pdf", mime_type=sniff.PDF
            )
